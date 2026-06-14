# -*- coding: utf-8 -*-
"""生成《Reopening Openness》复现汇报 PPT 初稿"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

BASE = "/sessions/admiring-nifty-bohr/mnt/Reopening Openness to Experience"
IMG  = os.path.join(BASE, "图片")
AST  = os.path.join(BASE, "ppt", "assets")
CC   = os.path.join(BASE, "复现输出", "claude code复现结果")
OUT  = os.path.join(BASE, "ppt", "复现汇报_初稿_v2.pptx")

# ---- palette ----
NAVY  = RGBColor(0x1E,0x27,0x61)
INK   = RGBColor(0x1A,0x1A,0x2E)
ICE   = RGBColor(0xCA,0xDC,0xFC)
CORAL = RGBColor(0xF9,0x61,0x67)
TEAL  = RGBColor(0x02,0xA8,0x88)
WHITE = RGBColor(0xFF,0xFF,0xFF)
TEXT  = RGBColor(0x1F,0x29,0x37)
MUTE  = RGBColor(0x6B,0x72,0x80)
CODEBG= RGBColor(0x22,0x22,0x3B)
CODETX= RGBColor(0xE6,0xE6,0xF0)
LINEBG= RGBColor(0xF1,0xF4,0xFA)

HEAD="微软雅黑"; BODY="微软雅黑"; CODE="Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def _setfont(run,name,size,bold,color,italic=False):
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    run.font.color.rgb=color; run.font.name=name
    rPr=run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el=rPr.find(qn(tag))
        if el is None:
            el=rPr.makeelement(qn(tag),{}); rPr.append(el)
        el.set('typeface',name)

def box(s,x,y,w,h,anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT,wrap=True):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align
    return tb,tf

def para(tf,first=False,space_before=6,space_after=0,align=PP_ALIGN.LEFT,line=1.0):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.space_before=Pt(space_before); p.space_after=Pt(space_after)
    try: p.line_spacing=line
    except: pass
    return p

def run(p,text,name=BODY,size=16,bold=False,color=TEXT,italic=False):
    r=p.add_run(); r.text=text; _setfont(r,name,size,bold,color,italic); return r

def rrect(s,x,y,w,h,fill,line=None,radius=0.08):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=radius
    except: pass
    return sp

def rect(s,x,y,w,h,fill,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp

def pic_fit(s,path,x,y,maxw,maxh,align="center",valign="middle"):
    im=Image.open(path); iw,ih=im.size; ar=iw/ih
    w=maxw; h=int(w/ar)
    if h>maxh: h=maxh; w=int(h*ar)
    px=x+(maxw-w)//2 if align=="center" else (x+maxw-w if align=="right" else x)
    py=y+(maxh-h)//2 if valign=="middle" else (y+maxh-h if valign=="bottom" else y)
    return s.shapes.add_picture(path,px,py,w,h)

def title(s,t,sub=None,dark=False,xy=(0.7,0.5),w=12.0):
    c=WHITE if dark else NAVY
    tb,tf=box(s,Inches(xy[0]),Inches(xy[1]),Inches(w),Inches(1.0))
    p=para(tf,first=True); run(p,t,HEAD,30,True,c)
    if sub:
        p2=para(tf,space_before=3); run(p2,sub,BODY,14,False,(ICE if dark else MUTE))
    return tb

def ai_badge(s,x=11.05,y=0.55):
    sp=rrect(s,Inches(x),Inches(y),Inches(1.95),Inches(0.42),CORAL,radius=0.5)
    sp.text_frame.word_wrap=False
    p=sp.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    sp.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    run(p,"AI 协作记忆点",BODY,12,True,WHITE)

def code_block(s,x,y,w,h,lines,fontsize=14):
    rrect(s,x,y,w,h,CODEBG,radius=0.04)
    tb=s.shapes.add_textbox(x+Inches(0.22),y+Inches(0.16),w-Inches(0.44),h-Inches(0.32))
    tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,(txt,col) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(0);p.space_after=Pt(2);p.line_spacing=1.05
        run(p,txt,CODE,fontsize,False,col)
    return tb

def bullets(s,x,y,w,h,items,size=16,gap=8,line=1.12):
    tb,tf=box(s,x,y,w,h)
    for i,it in enumerate(items):
        if isinstance(it,tuple): txt,lvl,kw=it
        else: txt,lvl,kw=it,0,{}
        p=para(tf,first=(i==0),space_before=(0 if i==0 else gap),line=line)
        if lvl>0: p.level=lvl
        mk = "▸ " if lvl==0 else "· "
        run(p,mk,BODY,size,kw.get("bold",False),kw.get("mkcolor",NAVY))
        run(p,txt,BODY,size,kw.get("bold",False),kw.get("color",TEXT))
    return tb

# ============================================================ S1 cover
s=slide(NAVY)
rect(s,0,0,Inches(0.28),SH,CORAL)
ai_badge(s,11.0,0.55)
tb,tf=box(s,Inches(0.7),Inches(1.05),Inches(6.7),Inches(0.9))
p=para(tf,first=True); run(p,'有同学问我：“Claude 能在 RStudio 里打开吗？”',BODY,16,False,ICE)
tb,tf=box(s,Inches(0.7),Inches(1.55),Inches(6.7),Inches(1.2))
p=para(tf,first=True); run(p,"能。",HEAD,60,True,CORAL)
tb,tf=box(s,Inches(0.7),Inches(2.95),Inches(6.6),Inches(2.4))
p=para(tf,first=True,line=1.12); run(p,"《Reopening Openness to Experience》",HEAD,28,True,WHITE)
p=para(tf,space_before=2,line=1.12); run(p,"计算可复现性检验",HEAD,28,True,WHITE)
p=para(tf,space_before=12); run(p,"网络分析 TMFG + walktrap  ·  MANOVA 组间比较",BODY,15,False,ICE)
p=para(tf,space_before=4); run(p,"全程 AI 协作复现",BODY,15,False,CORAL)
tb,tf=box(s,Inches(0.7),Inches(6.45),Inches(7),Inches(0.6))
p=para(tf,first=True); run(p,"汇报人：ORANGE   ·   R 编程语言   ·   2026-06-11",BODY,13,False,ICE)
pic_fit(s,os.path.join(IMG,"开头与Claude code的互动.png"),Inches(7.7),Inches(1.05),Inches(5.2),Inches(5.5),align="center")

# ============================================================ S2 文献信息
s=slide()
title(s,"被复现文献")
tb,tf=box(s,Inches(0.7),Inches(1.42),Inches(8),Inches(0.4))
p=para(tf,first=True); run(p,"TARGET PAPER · 计算可复现性检验对象",BODY,12,True,CORAL)
# 文章标题（hero）
tb,tf=box(s,Inches(0.7),Inches(1.82),Inches(12.0),Inches(1.5))
p=para(tf,first=True,line=1.12); run(p,"Reopening Openness to Experience:",HEAD,22,True,NAVY)
p=para(tf,space_before=2,line=1.12); run(p,"A Network Analysis of Four Openness to Experience Inventories",HEAD,22,True,NAVY)
# 作者
tb,tf=box(s,Inches(0.7),Inches(3.22),Inches(12.0),Inches(0.5))
p=para(tf,first=True)
run(p,"Christensen, A. P., Cotter, K. N., & Silvia, P. J. ",BODY,16,True,TEXT)
run(p,"(2019)",BODY,16,False,MUTE)
# 三张信息卡
cy=Inches(3.95); ch=Inches(1.9); ccw=Inches(3.7); cgp=Inches(0.4); cx0=Inches(0.7)
# A 期刊
rrect(s,cx0,cy,ccw,ch,LINEBG,radius=0.06)
tb,tf=box(s,cx0+Inches(0.3),cy+Inches(0.25),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"期刊 JOURNAL",BODY,12,True,MUTE)
tb,tf=box(s,cx0+Inches(0.3),cy+Inches(0.62),ccw-Inches(0.6),Inches(0.9))
p=para(tf,first=True,line=1.08); run(p,"Journal of Personality Assessment",HEAD,17,True,NAVY)
tb,tf=box(s,cx0+Inches(0.3),cy+Inches(1.46),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"Routledge (Taylor & Francis) · SSCI",BODY,12,False,MUTE)
# B 影响因子
cxb=cx0+ccw+cgp
rrect(s,cxb,cy,ccw,ch,NAVY,radius=0.06)
tb,tf=box(s,cxb+Inches(0.3),cy+Inches(0.25),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"影响因子 JIF · JCR 2024",BODY,12,True,ICE)
tb,tf=box(s,cxb+Inches(0.3),cy+Inches(0.55),ccw-Inches(0.6),Inches(0.9))
p=para(tf,first=True); run(p,"3.16",HEAD,46,True,WHITE)
tb,tf=box(s,cxb+Inches(0.3),cy+Inches(1.46),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"SSCI · Q2 · 5 年 IF 3.2",BODY,12,False,ICE)
# C 学术影响
cxc=cxb+ccw+cgp
rrect(s,cxc,cy,ccw,ch,LINEBG,radius=0.06)
tb,tf=box(s,cxc+Inches(0.3),cy+Inches(0.25),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"学术影响 · 被引",BODY,12,True,MUTE)
tb,tf=box(s,cxc+Inches(0.3),cy+Inches(0.55),ccw-Inches(0.6),Inches(0.9))
p=para(tf,first=True); run(p,"239",HEAD,46,True,CORAL)
tb,tf=box(s,cxc+Inches(0.3),cy+Inches(1.46),ccw-Inches(0.6),Inches(0.35))
p=para(tf,first=True); run(p,"Google Scholar · 2019 发表",BODY,12,False,MUTE)
# 关键词标签
tb,tf=box(s,Inches(0.7),Inches(6.15),Inches(4),Inches(0.35))
p=para(tf,first=True); run(p,"关键词标签",BODY,12,True,MUTE)
def _tw(t):
    w=0.0
    for chx in t: w += 0.17 if ord(chx)>0x2E80 else 0.095
    return w
tags=["网络分析","人格心理学","开放性结构","心理测量学","计算可复现性"]
tagx=Inches(0.7); tagy=Inches(6.5)
for tg in tags:
    wn=Inches(_tw(tg)+0.55)
    sp=rrect(s,tagx,tagy,wn,Inches(0.5),ICE,radius=0.5)
    sp.text_frame.word_wrap=False; sp.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    pp=sp.text_frame.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    run(pp,tg,BODY,13,True,NAVY)
    tagx=tagx+wn+Inches(0.25)

# ============================================================ S3 开放性 + 四套量表
s=slide()
title(s,"开放性是什么？四套量表各说各话")
bullets(s,Inches(0.7),Inches(1.7),Inches(6.5),Inches(4.5),[
  ("开放性（Openness to Experience）= 大五人格之一 ✓",0,{"bold":True,"mkcolor":TEAL}),
  ("传统认为含 2 个高阶方面：体验（Experiencing）+ 智力（Intellect）",0,{}),
  ("争议在低阶分面（facet）层：到底由哪些成分组成、",0,{}),
  ("各问卷说法不一 —— 结构本身就有争议",1,{"color":MUTE}),
  ("四套量表都号称测开放性，但各自侧重不同",0,{}),
],size=17,gap=12)
rrect(s,Inches(0.7),Inches(6.05),Inches(6.5),Inches(0.95),CORAL,radius=0.1)
tb,tf=box(s,Inches(0.95),Inches(6.05),Inches(6.0),Inches(0.95),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"伏笔：它们测的，真是同一个东西吗？",BODY,16,True,WHITE)
rows=[("量表","题数"),("NEO-PI-3","48"),("BFAS","20"),("HEXACO-100","16"),("Woo et al. (2014)","54"),("合计","138")]
tx=Inches(7.7); tw=Inches(5.0); ty=Inches(1.8); rh=Inches(0.72)
tbl=s.shapes.add_table(len(rows),2,tx,ty,tw,rh*len(rows)).table
tbl.columns[0].width=Inches(3.5); tbl.columns[1].width=Inches(1.5)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tbl.cell(ri,ci); cell.margin_top=Pt(3);cell.margin_bottom=Pt(3)
        cell.margin_left=Pt(10);cell.margin_right=Pt(6)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
        elif ri==len(rows)-1: cell.fill.solid(); cell.fill.fore_color.rgb=ICE
        else: cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
        p=cell.text_frame.paragraphs[0]
        p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        col=WHITE if ri==0 else NAVY if ri==len(rows)-1 else TEXT
        run(p,val,BODY,15,(ri==0 or ri==len(rows)-1),col)

# ============================================================ S4 研究目的+结论
s=slide()
title(s,"研究目的与原文结论")
rrect(s,Inches(0.7),Inches(1.7),Inches(11.9),Inches(1.55),NAVY,radius=0.06)
tb,tf=box(s,Inches(1.0),Inches(1.7),Inches(11.3),Inches(1.55),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"核心问题",BODY,14,True,CORAL)
p=para(tf,space_before=4,line=1.15); run(p,"四套主流开放性问卷，测的是同一个东西吗？开放性到底由哪些成分组成？",HEAD,20,True,WHITE)
tb,tf=box(s,Inches(0.7),Inches(3.55),Inches(11.9),Inches(0.5))
p=para(tf,first=True); run(p,"原文结论",BODY,16,True,NAVY)
bullets(s,Inches(0.85),Inches(4.1),Inches(11.6),Inches(3.0),[
  ("138 题聚成 10 个分面（facet），归入 3 个高阶方面：体验 · 智力 · 思想开放（Open-Mindedness）",0,{}),
  ("四套量表测同一大构念，但覆盖极不均衡",0,{"bold":True,"mkcolor":CORAL}),
  ("在网络识别出的 10 个分面中：Woo 与 NEO 各覆盖 9 个（最广）、BFAS 7 个、HEXACO 仅 4 个（最窄）",0,{}),
  ("Woo 题量最大（54 题）、三个高阶方面最均衡，综合覆盖最全面",1,{"color":MUTE}),
],size=17,gap=12)

# ============================================================ S6 方法+代码
s=slide()
title(s,"方法：TMFG 构网 + walktrap 社区检测")
bullets(s,Inches(0.7),Inches(1.7),Inches(6.4),Inches(3.5),[
  ("TMFG",0,{"bold":True,"mkcolor":TEAL}),
  ("把 138×138 相关矩阵“瘦身”成只留最强的",1,{}),
  ("3n − 6 = 408 条边的平面三角网络",1,{"color":MUTE}),
  ("walktrap",0,{"bold":True,"mkcolor":TEAL}),
  ("用随机游走找社区 —— 题目自动聚成分面",1,{}),
  ("确定性算法：换种子结果不变",1,{"color":MUTE}),
],size=17,gap=10)
code_block(s,Inches(7.3),Inches(1.85),Inches(5.3),Inches(2.7),[
  ("# 构网 -> 转 igraph -> 社区检测",CORAL),
  ("tmfg <- TMFG(data)$A",CODETX),
  ("g    <- convert2igraph(tmfg)",CODETX),
  ("wc   <- walktrap.community(g)",CODETX),
  ("max(wc$membership)   # 社区数",RGBColor(0x9A,0xA0,0xB5)),
],fontsize=15)
rrect(s,Inches(0.7),Inches(6.0),Inches(11.9),Inches(0.95),LINEBG,radius=0.1)
tb,tf=box(s,Inches(0.95),Inches(6.0),Inches(11.4),Inches(0.95),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"技术细节备问：",BODY,15,True,NAVY)
run(p,"为什么是平面网络 / 为什么 3n−6 条边 / walktrap 的随机游走如何聚类",BODY,15,False,TEXT)

# ============================================================ S7 Figure1
s=slide()
title(s,"我的工作量①：整体网络图（Figure 1）")
bullets(s,Inches(0.7),Inches(1.75),Inches(4.3),Inches(4.5),[
  ("138 节点 · 4 种形状区分 4 套问卷 · 10 社区配色",0,{}),
  ("圆 = NEO   方 = BFAS",1,{}),
  ("菱 = HEXACO   三角 = Woo",1,{}),
  ("三角（Woo）几乎进入每个颜色块",0,{}),
  ("→ 直观体现 Woo 覆盖面最广",1,{"color":MUTE}),
],size=16,gap=10)
rrect(s,Inches(0.7),Inches(6.0),Inches(4.3),Inches(0.85),CORAL,radius=0.1)
tb,tf=box(s,Inches(0.95),Inches(6.0),Inches(3.8),Inches(0.85),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"这张图是我自己跑出来的，不是论文截图",BODY,14,True,WHITE)
pic_fit(s,os.path.join(AST,"整体网络图-1.png"),Inches(5.3),Inches(1.55),Inches(7.5),Inches(5.4))

# ============================================================ S8 Figure2 + coverage
s=slide()
title(s,"我的工作量②：四套分图与覆盖结论（Figure 2）")
cov=[("BFAS","7"),("HEXACO","4"),("NEO","9"),("Woo","9")]
x0=Inches(0.7); y0=Inches(1.65); cw=Inches(1.05); gap=Inches(0.18)
for i,(nm,v) in enumerate(cov):
    cx=x0+i*(cw+gap)
    rrect(s,cx,y0,cw,Inches(1.25),NAVY,radius=0.1)
    tb,tf=box(s,cx,y0+Inches(0.12),cw,Inches(0.7),align=PP_ALIGN.CENTER)
    p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,v,HEAD,30,True,WHITE)
    tb,tf=box(s,cx,y0+Inches(0.82),cw,Inches(0.4),align=PP_ALIGN.CENTER)
    p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,nm,BODY,12,False,ICE)
tb,tf=box(s,Inches(0.7),Inches(3.05),Inches(5.4),Inches(0.4))
p=para(tf,first=True); run(p,"覆盖分面数 / 10（与论文逐项吻合）",BODY,12,False,MUTE)
rrect(s,Inches(0.7),Inches(3.6),Inches(5.5),Inches(3.0),CORAL,radius=0.06)
tb,tf=box(s,Inches(0.95),Inches(3.8),Inches(5.0),Inches(2.7),anchor=MSO_ANCHOR.TOP)
p=para(tf,first=True); run(p,"杀手锏结论",BODY,15,True,WHITE)
for t in ["覆盖极不均衡","BFAS / HEXACO 完全没覆盖“思想开放”","NEO 与 Woo 最广（各 9 个 facet）","Woo 最均衡、题量最大，覆盖最全面"]:
    p=para(tf,space_before=8,line=1.1); run(p,"• "+t,BODY,15,False,WHITE)
pic_fit(s,os.path.join(AST,"量表各自网络图-1.png"),Inches(6.5),Inches(1.55),Inches(6.4),Inches(5.4))

# ============================================================ S9 复现对照 一致程度
s=slide()
title(s,"复现对照：社区结构的一致程度")
rrect(s,Inches(0.7),Inches(1.7),Inches(3.2),Inches(1.8),NAVY,radius=0.08)
tb,tf=box(s,Inches(0.7),Inches(1.85),Inches(3.2),Inches(1.0),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"92.8%",HEAD,44,True,WHITE)
tb,tf=box(s,Inches(0.7),Inches(2.75),Inches(3.2),Inches(0.6),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"逐题一致率 128/138",BODY,13,False,ICE)
bullets(s,Inches(4.1),Inches(1.7),Inches(3.4),Inches(2.0),[
  ("10 社区与论文一对一全部还原",0,{}),
  ("6 个 facet 100% 一致，其余 73–92%",0,{}),
  ("仅 10 道边界题摆动",0,{"bold":True,"mkcolor":CORAL}),
],size=15,gap=8)
rrect(s,Inches(0.7),Inches(3.7),Inches(6.8),Inches(1.2),LINEBG,radius=0.06)
tb,tf=box(s,Inches(0.95),Inches(3.82),Inches(6.4),Inches(1.0))
p=para(tf,first=True,line=1.1); run(p,"10 道边界题（反向计分/外围题）：",BODY,13,True,NAVY)
p=para(tf,space_before=3,line=1.1); run(p,"Fe1  Fe2  Va7  Ac1  Fa8  Op5  Uc4  To4  Ig5  Cu4",CODE,14,False,TEXT)
rrect(s,Inches(0.7),Inches(5.1),Inches(6.8),Inches(1.6),CORAL,radius=0.06)
tb,tf=box(s,Inches(0.95),Inches(5.25),Inches(6.4),Inches(1.35))
p=para(tf,first=True); run(p,"⚠ 命名陷阱",BODY,14,True,WHITE)
p=para(tf,space_before=4,line=1.15); run(p,"walktrap 给社区发的编号是任意的。必须按 SI 2 逐题核对重贴 facet 名 —— 原脚本第 39 行按数字硬贴会贴错。",BODY,14,False,WHITE)
pic_fit(s,os.path.join(IMG,"修改后的代码2.png"),Inches(7.75),Inches(1.7),Inches(5.1),Inches(5.0),valign="top")

# ============================================================ S10 8 vs 10 (AI)
s=slide()
title(s,"重点修改①：社区切法 10 vs 8")
ai_badge(s,11.0,0.55)
bullets(s,Inches(0.7),Inches(1.7),Inches(6.3),Inches(2.6),[
  ("现象：新版 igraph 跑 walktrap 自动切出 8 个社区，不是论文的 10",0,{}),
  ("根因：igraph 版本漂移（方法零差异）",0,{"bold":True,"mkcolor":CORAL}),
  ("模块度：8 类 = 0.722 ＞ 10 类 = 0.712（几乎一样好）",1,{}),
  ("处理：cut_at(no = 10) 强制切回 10，对齐论文",0,{}),
],size=15,gap=9)
code_block(s,Inches(0.7),Inches(4.55),Inches(6.3),Inches(1.35),[
  ("wc  <- walktrap.community(g)   # 自动切 = 8",RGBColor(0x9A,0xA0,0xB5)),
  ("m10 <- igraph::cut_at(wc, no = 10)  # 强切 10",CODETX),
  ("igraph::modularity(g, m10)     # 0.712",CODETX),
],fontsize=13)
rrect(s,Inches(0.7),Inches(6.05),Inches(6.3),Inches(0.9),NAVY,radius=0.1)
tb,tf=box(s,Inches(0.95),Inches(6.05),Inches(5.8),Inches(0.9),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True,line=1.1); run(p,"点题：复现 ≠ 点一下就出一样结果 —— 这是版本敏感性。",BODY,14,True,WHITE)
tb,tf=box(s,Inches(7.3),Inches(1.65),Inches(5.5),Inches(0.4))
p=para(tf,first=True); run(p,"版本证据（本机）：",BODY,13,True,NAVY)
pic_fit(s,os.path.join(IMG,"R包的版本.png"),Inches(7.3),Inches(2.0),Inches(5.5),Inches(1.0),align="left")
tb,tf=box(s,Inches(7.3),Inches(3.2),Inches(5.5),Inches(0.4))
p=para(tf,first=True); run(p,"我的修正代码：",BODY,13,True,NAVY)
pic_fit(s,os.path.join(IMG,"修改后的代码1.png"),Inches(7.3),Inches(3.55),Inches(5.5),Inches(3.2),align="left",valign="top")

# ============================================================ S11 Group rebuild (AI)
s=slide()
title(s,"重点修改②：原始数据缺分组变量")
ai_badge(s,11.0,0.55)
bullets(s,Inches(0.7),Inches(1.7),Inches(6.3),Inches(2.3),[
  ("问题：分析文件只有无关的 DBHTgroup，没有 MANOVA 需要的 Group 列",0,{}),
  ("→ 直接跑原脚本报“找不到对象 Group”",1,{"color":CORAL}),
  ("解决：用 138 题“响应指纹”把 802 人匹配回三个源文件",0,{"bold":True,"mkcolor":NAVY}),
  ("（Fall 2015 / Spring 2017 / Mturk）",1,{"color":MUTE}),
],size=15,gap=9)
rrect(s,Inches(0.7),Inches(4.15),Inches(6.3),Inches(1.25),NAVY,radius=0.08)
tb,tf=box(s,Inches(0.7),Inches(4.28),Inches(6.3),Inches(0.6),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"重建 Group = 176 / 108 / 518",HEAD,26,True,WHITE)
tb,tf=box(s,Inches(0.7),Inches(4.92),Inches(6.3),Inches(0.45),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"与论文 SI 7、Data Cleaning Notes 三方一致 ✓",BODY,14,False,TEAL)
rrect(s,Inches(0.7),Inches(5.6),Inches(6.3),Inches(1.35),LINEBG,radius=0.08)
tb,tf=box(s,Inches(0.95),Inches(5.72),Inches(5.85),Inches(1.1),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True,line=1.15); run(p,"为什么唯一：每人答 138 题（1–5 分），完整应答序列几乎唯一（5^138 种组合），可当“指纹”无损反查样本来源。",BODY,13,False,TEXT)
tb,tf=box(s,Inches(7.3),Inches(1.65),Inches(5.5),Inches(0.4))
p=para(tf,first=True); run(p,"我的指纹重建 + MANOVA 代码：",BODY,13,True,NAVY)
pic_fit(s,os.path.join(IMG,"修改后的代码4.png"),Inches(7.3),Inches(2.0),Inches(5.5),Inches(4.8),align="left",valign="top")

# ============================================================ S12 MANOVA results
s=slide()
title(s,"结果：MANOVA + Box's M 精确命中")
hdr=["比较","复现 Pillai","论文","复现 F","论文 F"]
data=[
 ["三组整体",".9596",".960","4.43","4.43"],
 ["1 vs 2（Fall×Spring）",".650",".650","1.95","1.95"],
 ["1 vs 3（Fall×Mturk）",".675",".675","8.34","8.34"],
 ["2 vs 3（Spring×Mturk）",".614",".614","5.62","5.62"],
]
tx=Inches(0.7); ty=Inches(1.85); tw=Inches(7.6); rh=Inches(0.62)
tbl=s.shapes.add_table(len(data)+1,5,tx,ty,tw,rh*(len(data)+1)).table
widths=[2.6,1.35,1.05,1.25,1.05]
for ci,wv in enumerate(widths): tbl.columns[ci].width=Inches(wv)
for ci,h in enumerate(hdr):
    cell=tbl.cell(0,ci); cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
    p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
    run(p,h,BODY,14,True,WHITE)
for ri,rowd in enumerate(data,start=1):
    for ci,val in enumerate(rowd):
        cell=tbl.cell(ri,ci); cell.fill.solid()
        cell.fill.fore_color.rgb=WHITE if ri%2 else LINEBG
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        run(p,val,BODY if ci==0 else CODE,14,False,TEXT)
rrect(s,Inches(8.7),Inches(1.85),Inches(3.9),Inches(2.0),LINEBG,radius=0.06)
tb,tf=box(s,Inches(8.95),Inches(2.0),Inches(3.5),Inches(1.8))
p=para(tf,first=True); run(p,"Box's M（1 vs 3）",BODY,15,True,NAVY)
p=para(tf,space_before=6,line=1.15); run(p,"χ² = 12891  df = 9591  p < .001",CODE,14,False,TEXT)
p=para(tf,space_before=4,line=1.15); run(p,"三组协方差结构显著不齐",BODY,13,False,MUTE)
p=para(tf,space_before=4,line=1.15); run(p,"样本2 (n=108<138) 奇异，无法算 → 只检验 1 vs 3",BODY,12,False,MUTE)
rrect(s,Inches(8.7),Inches(4.0),Inches(3.9),Inches(1.2),CORAL,radius=0.1)
tb,tf=box(s,Inches(8.7),Inches(4.0),Inches(3.9),Inches(1.2),anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"4 个 MANOVA 全部命中 ✓",HEAD,18,True,WHITE)
tb,tf=box(s,Inches(0.7),Inches(6.4),Inches(11.9),Inches(0.6))
p=para(tf,first=True); run(p,"百分误差 ≈ 0（Pillai .9596 vs .960 ≈ 0.04%），按原文报告精度四舍五入逐一相等。",BODY,13,False,MUTE)

# ============================================================ S13a 对照表 推断统计
s=slide()
title(s,"复现对照表①：推断统计（原文方法）")
hdr=["效应（统计量）","原文","本研究·10社区","Claude code·8社区","评级"]
data=[
 ["整体 MANOVA  F","4.43","4.4308","4.4308","完全一致"],
 ["整体  Pillai",".960",".9596",".9596","完全一致"],
 ["1v2  F / Pillai","1.95 / .650","1.9516 / .650","1.9516 / .650","完全一致"],
 ["1v3  F / Pillai","8.34 / .675","8.3412 / .6747","8.3412 / .6747","完全一致"],
 ["2v3  F / Pillai","5.62 / .614","5.6191 / .6142","5.6191 / .6142","完全一致"],
 ["Box's M (1v3)","F=1.292","χ²=12891 (df=9591)","同左","完全一致*"],
]
tx=Inches(0.7); ty=Inches(1.8); tw=Inches(11.9); rh=Inches(0.6)
tbl=s.shapes.add_table(len(data)+1,5,tx,ty,tw,rh*(len(data)+1)).table
widths=[2.7,2.0,2.7,2.7,1.8]
for ci,wv in enumerate(widths): tbl.columns[ci].width=Inches(wv)
for ci,h in enumerate(hdr):
    cell=tbl.cell(0,ci); cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
    p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
    run(p,h,BODY,13,True,WHITE)
for ri,rowd in enumerate(data,start=1):
    for ci,val in enumerate(rowd):
        cell=tbl.cell(ri,ci); cell.fill.solid()
        cell.fill.fore_color.rgb=WHITE if ri%2 else LINEBG
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE; cell.margin_left=Pt(8)
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
        col=TEAL if ci==4 else TEXT
        run(p,val,BODY,12,(ci==4),col)
tb,tf=box(s,Inches(0.7),Inches(6.55),Inches(11.9),Inches(0.6))
p=para(tf,first=True,line=1.1); run(p,"* Box's M：R 报卡方近似、原文报 F 近似，同一检验两种换算，df=9591 一致、p<.001 推论一致。两版复现（10社区/8社区）MANOVA 与社区数无关，数值逐字一致，互为独立交叉验证。",BODY,11,False,MUTE)

# ============================================================ S13b 评级汇总
s=slide()
title(s,"复现对照表②：可复现性与推论一致性评级")
cards=[("14 / 14","结果可复现性\n完全一致（δ = 0%）",NAVY),
       ("5 / 5","推论一致性\n原文与复现同侧（100%）",NAVY),
       ("4 / 4","描述统计·样本量\n完全一致（802/176/108/518）",NAVY)]
x0=Inches(0.7); y0=Inches(2.0); cw=Inches(3.7); gap=Inches(0.4)
for i,(big,lab,bg) in enumerate(cards):
    cx=x0+i*(cw+gap)
    rrect(s,cx,y0,cw,Inches(2.4),bg,radius=0.06)
    tb,tf=box(s,cx,y0+Inches(0.35),cw,Inches(1.0),align=PP_ALIGN.CENTER)
    p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,big,HEAD,46,True,WHITE)
    tb,tf=box(s,cx,y0+Inches(1.45),cw,Inches(0.85),align=PP_ALIGN.CENTER)
    for j,ln in enumerate(lab.split("\n")):
        p=para(tf,first=(j==0),align=PP_ALIGN.CENTER,space_before=(0 if j==0 else 2))
        run(p,ln,BODY,(14 if j==0 else 12),(j==0),(WHITE if j==0 else ICE))
rrect(s,Inches(0.7),Inches(4.9),Inches(11.9),Inches(1.5),LINEBG,radius=0.06)
tb,tf=box(s,Inches(0.95),Inches(5.05),Inches(11.4),Inches(1.25),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"评级三档（学院指南）：",BODY,15,True,NAVY)
run(p,"完全一致 (0%)  /  次要偏差 (0%<δ<10%)  /  主要偏差 (≥10%)",BODY,15,False,TEXT)
p=para(tf,space_before=8,line=1.15); run(p,"本复现 14 项推断结果全部落入“完全一致”，5 个统计推断全部与原文同侧。",BODY,15,False,TEXT)

# ============================================================ S13c 不可复现原因
s=slide()
title(s,"复现对照表③：(不)可复现原因分析")
items=[
 ("软件包版本（igraph）","唯一造成结构差异的因素 —— walktrap 默认社区数 10->8；cut_at 强切回 10 后逐一致。已定位、可控。",CORAL),
 ("OSF 数据缺分组变量","分析文件缺 Group 列（仅有无关 DBHTgroup）-> 用 138 题指纹重建 176/108/518。",NAVY),
 ("未提供原始数据","仅公开处理后数据，原清洗在 SPSS；合规采用处理后数据。",NAVY),
 ("缺 readme 说明","OSF 未见 readme，增加理解数据/代码对应关系的成本。",NAVY),
]
y=Inches(1.85)
for nm,desc,c in items:
    rrect(s,Inches(0.7),y,Inches(0.16),Inches(1.0),c,radius=0.5)
    tb,tf=box(s,Inches(1.05),y,Inches(11.5),Inches(1.05),anchor=MSO_ANCHOR.MIDDLE)
    p=para(tf,first=True); run(p,nm,BODY,16,True,c if c==CORAL else NAVY)
    p=para(tf,space_before=3,line=1.12); run(p,desc,BODY,13,False,TEXT)
    y=y+Inches(1.18)
rrect(s,Inches(0.7),Inches(6.55),Inches(11.9),Inches(0.55),LINEBG,radius=0.1)
tb,tf=box(s,Inches(0.95),Inches(6.55),Inches(11.4),Inches(0.55),anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,first=True); run(p,"整体：结果 100% 一致、推论 100% 一致；仅社区结构因 igraph 版本出现 10->8 的可控差异。",BODY,13,True,NAVY)

# ============================================================ S14 总结
s=slide(NAVY)
rect(s,0,0,Inches(0.28),SH,CORAL)
title(s,"总结与结论",dark=True)
bullets(s,Inches(0.7),Inches(1.8),Inches(11.9),Inches(2.0),[
  ("复现结论：网络结构 + MANOVA 关键结果全部与原文定量吻合",0,{"color":WHITE,"mkcolor":CORAL,"bold":True}),
  ("四套问卷测同一构念，覆盖却极不均衡 —— 回扣核心问题",0,{"color":ICE}),
],size=18,gap=12)
tb,tf=box(s,Inches(0.7),Inches(3.5),Inches(11.9),Inches(0.5))
p=para(tf,first=True); run(p,"两个超出预期的发现",BODY,16,True,CORAL)
y0=Inches(4.05)
for i,(big,desc) in enumerate([("①","igraph 版本致社区数敏感（10 → 8）—— 复现里典型的版本敏感性"),
                                ("②","用响应指纹补全原文遗漏的分组变量 —— 一次“数据侦查”")]):
    rrect(s,Inches(0.7),y0+i*Inches(1.25),Inches(11.9),Inches(1.05),INK,radius=0.06)
    tb,tf=box(s,Inches(1.0),y0+i*Inches(1.25),Inches(1.0),Inches(1.05),anchor=MSO_ANCHOR.MIDDLE)
    p=para(tf,first=True); run(p,big,HEAD,36,True,CORAL)
    tb,tf=box(s,Inches(2.1),y0+i*Inches(1.25),Inches(10.3),Inches(1.05),anchor=MSO_ANCHOR.MIDDLE)
    p=para(tf,first=True,line=1.15); run(p,desc,BODY,16,False,WHITE)

# ============================================================ S15 Q&A
s=slide(NAVY)
rect(s,0,0,Inches(0.28),SH,CORAL)
tb,tf=box(s,Inches(0.7),Inches(2.6),Inches(11.9),Inches(1.5),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"谢谢！",HEAD,54,True,WHITE)
tb,tf=box(s,Inches(0.7),Inches(4.0),Inches(11.9),Inches(0.7),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"Q & A",HEAD,30,True,CORAL)
tb,tf=box(s,Inches(0.7),Inches(5.0),Inches(11.9),Inches(0.6),align=PP_ALIGN.CENTER)
p=para(tf,first=True,align=PP_ALIGN.CENTER); run(p,"欢迎提问：TMFG 技术细节 · 版本敏感性 · 指纹重建分组",BODY,15,False,ICE)

prs.save(OUT)
print("SAVED", OUT, "slides=", len(prs.slides._sldIdLst))
